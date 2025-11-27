"""
Native PPTX converter using python-pptx and hybrid screenshot approach.

This module provides functionality to convert presentations to editable PPTX format
by mapping CSS Grid layouts to PowerPoint shapes and text boxes.
"""

import logging
import io
from typing import Optional, Tuple, List, Dict, Any
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image

from .base import BaseConverter

logger = logging.getLogger(__name__)


class NativePPTXConverter(BaseConverter):
    """Convert presentations to native editable PPTX format."""

    # PowerPoint slide dimensions (16:9)
    PPTX_WIDTH = Inches(10)
    PPTX_HEIGHT = Inches(5.625)
    
    # Grid System (32x18)
    GRID_COLS = 32
    GRID_ROWS = 18
    
    # Calculate grid unit size
    COL_WIDTH = PPTX_WIDTH / GRID_COLS
    ROW_HEIGHT = PPTX_HEIGHT / GRID_ROWS

    async def generate_pptx(
        self,
        presentation_id: str,
        slide_count: int,
        output_path: Optional[Path] = None,
        presentation_data: Optional[Dict[str, Any]] = None
    ) -> bytes:
        """
        Generate a native PPTX file from a presentation.

        Args:
            presentation_id: The UUID of the presentation
            slide_count: Number of slides (used for progress tracking)
            output_path: Optional file path to save the PPTX
            presentation_data: The full JSON data of the presentation (required for native generation)

        Returns:
            PPTX file as bytes
        """
        logger.info(f"Generating Native PPTX for: {presentation_id}")

        if not presentation_data:
            # TODO: Fetch presentation data if not provided
            # For now, we assume it's passed or we need to fetch it
            logger.warning("Presentation data not provided, fetching from API...")
            presentation_data = await self._fetch_presentation_data(presentation_id)

        # Create PowerPoint presentation
        prs = Presentation()
        prs.slide_width = self.PPTX_WIDTH
        prs.slide_height = self.PPTX_HEIGHT

        # Process each slide
        slides_data = presentation_data.get('slides', [])
        for idx, slide_data in enumerate(slides_data):
            logger.info(f"Processing slide {idx + 1}/{len(slides_data)}")
            
            # Create blank slide
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Apply background
            self._apply_background(slide, slide_data)
            
            # Render layout
            layout_type = slide_data.get('layout', 'L01')
            content = slide_data.get('content', {})
            
            if layout_type == 'L01':
                await self._render_L01(slide, content, idx, presentation_id)
            elif layout_type == 'L02':
                await self._render_L02(slide, content, idx, presentation_id)
            elif layout_type == 'L03':
                await self._render_L03(slide, content, idx, presentation_id)
            elif layout_type == 'L25':
                await self._render_L25(slide, content, idx, presentation_id)
            elif layout_type == 'L27':
                await self._render_L27(slide, content, idx, presentation_id)
            elif layout_type == 'L29':
                await self._render_L29(slide, content, idx, presentation_id)
            else:
                logger.warning(f"Unsupported layout {layout_type}, skipping content")

        # Save to bytes
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_bytes = pptx_buffer.getvalue()

        if output_path:
            output_path.write_bytes(pptx_bytes)
            logger.info(f"Native PPTX saved to: {output_path}")

        return pptx_bytes

    def _grid_to_inches(self, col_start: int, col_end: int, row_start: int, row_end: int) -> Tuple[Inches, Inches, Inches, Inches]:
        """
        Convert CSS Grid coordinates to PowerPoint position and size.
        
        Args:
            col_start: Starting column line (1-based)
            col_end: Ending column line (1-based)
            row_start: Starting row line (1-based)
            row_end: Ending row line (1-based)
            
        Returns:
            Tuple of (left, top, width, height) in Inches
        """
        # CSS Grid is 1-based, so subtract 1 for 0-based offset
        left = (col_start - 1) * self.COL_WIDTH
        top = (row_start - 1) * self.ROW_HEIGHT
        
        # Width/Height is the span
        width = (col_end - col_start) * self.COL_WIDTH
        height = (row_end - row_start) * self.ROW_HEIGHT
        
        return (left, top, width, height)

    def _apply_background(self, slide, slide_data):
        """Apply background color or image to slide."""
        bg_color = slide_data.get('background_color')
        # TODO: Handle background images
        
        if bg_color:
            background = slide.background
            fill = background.fill
            fill.solid()
            # Parse hex color (e.g., #FFFFFF)
            if bg_color.startswith('#'):
                r = int(bg_color[1:3], 16)
                g = int(bg_color[3:5], 16)
                b = int(bg_color[5:7], 16)
                fill.fore_color.rgb = RGBColor(r, g, b)

    async def _render_L01(self, slide, content, slide_index, presentation_id):
        """
        Render L01 Layout: Centered Chart/Diagram with text.
        
        L01 Grid Mapping:
        - Title: row 2/3, col 2/32
        - Subtitle: row 3/4, col 2/32
        - Chart: row 5/15, col 2/32 (Hybrid Capture)
        - Body: row 15/17, col 2/32
        - Footer: row 18/19
        """
        # 1. Title
        self._add_text_box(
            slide,
            text=content.get('slide_title', ''),
            grid=(2, 32, 2, 3),
            font_size=42,
            is_bold=True,
            color=RGBColor(31, 41, 55)  # #1f2937
        )
        
        # 2. Subtitle
        self._add_text_box(
            slide,
            text=content.get('element_1', ''),
            grid=(2, 32, 3, 4),
            font_size=24,
            color=RGBColor(107, 114, 128)  # #6b7280
        )
        
        # 3. Chart (Hybrid Capture)
        # We need to capture the element at .chart-container
        # Selector: [data-slide-index="{slide_index}"] .chart-container
        selector = f'[data-slide-index="{slide_index}"] .chart-container'
        chart_bytes = await self.capture_element_screenshot(presentation_id, slide_index, selector)
        
        if chart_bytes:
            left, top, width, height = self._grid_to_inches(2, 32, 5, 15)
            slide.shapes.add_picture(io.BytesIO(chart_bytes), left, top, width, height)
            
        # 4. Body Text
        self._add_text_box(
            slide,
            text=content.get('element_3', ''),
            grid=(2, 32, 15, 17),
            font_size=20,
            color=RGBColor(55, 65, 81)  # #374151
        )
        
        # 5. Footer (Presentation Name)
        if content.get('presentation_name'):
            self._add_text_box(
                slide,
                text=content.get('presentation_name'),
                grid=(2, 7, 18, 19),
                font_size=18,
                color=RGBColor(31, 41, 55)
            )

    async def _render_L02(self, slide, content, slide_index, presentation_id):
        """
        Render L02 Layout: Left Diagram with Text Box on Right.
        
        L02 Grid Mapping:
        - Title: row 2/3, col 2/32
        - Subtitle: row 3/4, col 2/32
        - Diagram (Left): row 5/17, col 2/23 (Hybrid Capture)
        - Text (Right): row 5/17, col 23/32
        """
        # 1. Title
        self._add_text_box(
            slide,
            text=content.get('slide_title', ''),
            grid=(2, 32, 2, 3),
            font_size=42,
            is_bold=True,
            color=RGBColor(31, 41, 55)
        )
        
        # 2. Subtitle
        self._add_text_box(
            slide,
            text=content.get('element_1', ''),
            grid=(2, 32, 3, 4),
            font_size=24,
            color=RGBColor(107, 114, 128)
        )
        
        # 3. Diagram (Left) - Hybrid Capture
        # Selector: [data-slide-index="{slide_index}"] .diagram-container
        selector = f'[data-slide-index="{slide_index}"] .diagram-container'
        diagram_bytes = await self.capture_element_screenshot(presentation_id, slide_index, selector)
        
        if diagram_bytes:
            left, top, width, height = self._grid_to_inches(2, 23, 5, 17)
            slide.shapes.add_picture(io.BytesIO(diagram_bytes), left, top, width, height)
            
        # 4. Text (Right)
        # Note: element_2 might be HTML, but we'll render as plain text for now
        text_content = content.get('element_2', '')
        # Simple HTML tag stripping if needed, or just dump it
        # For now, assuming it's readable text
        self._add_text_box(
            slide,
            text=text_content,
            grid=(23, 32, 5, 17),
            font_size=20,
            color=RGBColor(55, 65, 81)
        )
        
        # 5. Footer
        if content.get('presentation_name'):
            self._add_text_box(
                slide,
                text=content.get('presentation_name'),
                grid=(2, 7, 18, 19),
                font_size=18,
                color=RGBColor(31, 41, 55)
            )

    async def _render_L03(self, slide, content, slide_index, presentation_id):
        """
        Render L03 Layout: Two Charts in Columns with Text Below.
        """
        # 1. Title
        self._add_text_box(slide, content.get('slide_title', ''), (2, 32, 2, 3), 42, True, RGBColor(31, 41, 55))
        
        # 2. Subtitle
        self._add_text_box(slide, content.get('element_1', ''), (2, 32, 3, 4), 24, False, RGBColor(107, 114, 128))
        
        # 3. Left Chart (Hybrid)
        selector_left = f'[data-slide-index="{slide_index}"] [data-section-type="chart1"]'
        chart1_bytes = await self.capture_element_screenshot(presentation_id, slide_index, selector_left)
        if chart1_bytes:
            left, top, width, height = self._grid_to_inches(2, 16, 5, 14)
            slide.shapes.add_picture(io.BytesIO(chart1_bytes), left, top, width, height)
            
        # 4. Right Chart (Hybrid)
        selector_right = f'[data-slide-index="{slide_index}"] [data-section-type="chart2"]'
        chart2_bytes = await self.capture_element_screenshot(presentation_id, slide_index, selector_right)
        if chart2_bytes:
            left, top, width, height = self._grid_to_inches(17, 31, 5, 14)
            slide.shapes.add_picture(io.BytesIO(chart2_bytes), left, top, width, height)
            
        # 5. Left Body
        self._add_text_box(slide, content.get('element_3', ''), (2, 16, 14, 17), 20, False, RGBColor(55, 65, 81))
        
        # 6. Right Body
        self._add_text_box(slide, content.get('element_5', ''), (17, 31, 14, 17), 20, False, RGBColor(55, 65, 81))
        
        # 7. Footer
        if content.get('presentation_name'):
            self._add_text_box(slide, content.get('presentation_name'), (2, 7, 18, 19), 18, False, RGBColor(31, 41, 55))

    async def _render_L25(self, slide, content, slide_index, presentation_id):
        """
        Render L25 Layout: Main Content Shell (Hybrid Rich Content).
        """
        # 1. Title
        self._add_text_box(slide, content.get('slide_title', ''), (2, 32, 2, 3), 42, True, RGBColor(31, 41, 55))
        
        # 2. Subtitle
        subtitle = content.get('subtitle') or content.get('element_1')
        self._add_text_box(slide, subtitle, (2, 32, 3, 4), 24, False, RGBColor(107, 114, 128))
        
        # 3. Rich Content (Hybrid Capture)
        # Selector: [data-slide-index="{slide_index}"] .rich-content-area
        selector = f'[data-slide-index="{slide_index}"] .rich-content-area'
        content_bytes = await self.capture_element_screenshot(presentation_id, slide_index, selector)
        if content_bytes:
            left, top, width, height = self._grid_to_inches(2, 32, 5, 17)
            slide.shapes.add_picture(io.BytesIO(content_bytes), left, top, width, height)
            
        # 4. Footer
        if content.get('presentation_name'):
            self._add_text_box(slide, content.get('presentation_name'), (2, 7, 18, 19), 18, False, RGBColor(31, 41, 55))

    async def _render_L27(self, slide, content, slide_index, presentation_id):
        """
        Render L27 Layout: Image Left with Content Right.
        """
        # 1. Left Image (Hybrid Capture of container to handle bg/img logic)
        selector = f'[data-slide-index="{slide_index}"] .image-container'
        img_bytes = await self.capture_element_screenshot(presentation_id, slide_index, selector)
        if img_bytes:
            left, top, width, height = self._grid_to_inches(1, 12, 1, 19)
            slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width, height)
            
        # 2. Title
        self._add_text_box(slide, content.get('slide_title', ''), (13, 32, 2, 3), 42, True, RGBColor(31, 41, 55))
        
        # 3. Subtitle
        self._add_text_box(slide, content.get('element_1', ''), (13, 32, 3, 4), 24, False, RGBColor(107, 114, 128))
        
        # 4. Main Content
        self._add_text_box(slide, content.get('main_content', ''), (13, 32, 5, 17), 20, False, RGBColor(55, 65, 81))
        
        # 5. Footer
        if content.get('presentation_name'):
            self._add_text_box(slide, content.get('presentation_name'), (13, 18, 18, 19), 18, False, RGBColor(31, 41, 55))

    async def _render_L29(self, slide, content, slide_index, presentation_id):
        """
        Render L29 Layout: Hero Full-Bleed (Full Slide Hybrid Capture).
        """
        # Capture the entire hero content area which spans the whole slide
        selector = f'[data-slide-index="{slide_index}"] .hero-content-area'
        hero_bytes = await self.capture_element_screenshot(presentation_id, slide_index, selector)
        if hero_bytes:
            left, top, width, height = self._grid_to_inches(1, 33, 1, 19) # Full slide 1-32 cols, 1-18 rows
            slide.shapes.add_picture(io.BytesIO(hero_bytes), left, top, width, height)
            
        col_start, col_end, row_start, row_end = grid
        left, top, width, height = self._grid_to_inches(col_start, col_end, row_start, row_end)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = is_bold
        p.alignment = align
        
        if color:
            p.font.color.rgb = color

    async def _fetch_presentation_data(self, presentation_id: str) -> Dict[str, Any]:
        """Fetch presentation JSON data from API."""
        # This requires the base_url to be set correctly to the API server
        # We might need to adjust base_url logic if API is on a different port/path
        import aiohttp
        
        # Assuming API is at {base_url}/api/presentations/{id}
        # We need to handle the case where base_url is the viewer URL
        # For now, let's assume we can construct it.
        # If running in docker, we might need internal DNS.
        
        # HACK: For now, we'll try to fetch from the same base URL
        api_url = f"{self.base_url}/api/presentations/{presentation_id}"
        
        async with aiohttp.ClientSession() as session:
            async with session.get(api_url) as resp:
                if resp.status == 200:
                    return await resp.json()
                else:
                    logger.error(f"Failed to fetch presentation data: {resp.status}")
                    return {}
