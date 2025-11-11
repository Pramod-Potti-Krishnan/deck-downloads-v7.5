"""
PPTX converter using python-pptx and Playwright screenshots.

This module provides functionality to convert presentations to PowerPoint
format by capturing slides as screenshots and embedding them in PPTX.
"""

import io
import logging
from pathlib import Path
from typing import Optional
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from .base import BaseConverter

logger = logging.getLogger(__name__)


class PPTXConverter(BaseConverter):
    """Convert presentations to PPTX format using screenshots and python-pptx."""

    # PowerPoint slide dimensions for 16:9 aspect ratio
    PPTX_WIDTH = Inches(10)
    PPTX_HEIGHT = Inches(5.625)  # 10 * 9/16 for 16:9 aspect ratio

    async def generate_pptx(
        self,
        presentation_id: str,
        slide_count: int,
        output_path: Optional[Path] = None,
        aspect_ratio: str = "16:9",
        quality: str = "high"
    ) -> bytes:
        """
        Generate a PPTX file from a presentation.

        This method captures each slide as a screenshot and embeds them
        in a PowerPoint presentation.

        Args:
            presentation_id: The UUID of the presentation
            slide_count: Number of slides in the presentation
            output_path: Optional file path to save the PPTX
            aspect_ratio: Aspect ratio - '16:9' or '4:3' (default: '16:9')
            quality: Quality setting - 'high', 'medium', or 'low'

        Returns:
            PPTX file as bytes

        Raises:
            ValueError: If presentation cannot be loaded or invalid parameters
            RuntimeError: If PPTX generation fails
        """
        logger.info(f"Generating PPTX for presentation: {presentation_id}")
        logger.info(f"Aspect ratio: {aspect_ratio}, Quality: {quality}, Slides: {slide_count}")

        # Validate parameters
        if aspect_ratio not in ["16:9", "4:3"]:
            raise ValueError(f"Invalid aspect ratio: {aspect_ratio}. Must be '16:9' or '4:3'")

        if slide_count <= 0:
            raise ValueError(f"Invalid slide count: {slide_count}")

        try:
            # Capture screenshots of all slides
            logger.info("Capturing slide screenshots...")
            screenshots = await self.capture_slide_screenshots(presentation_id, slide_count)

            if not screenshots:
                raise RuntimeError("No screenshots captured")

            logger.info(f"Captured {len(screenshots)} slide screenshots")

            # Create PowerPoint presentation
            prs = Presentation()

            # Set slide dimensions based on aspect ratio
            prs.slide_width, prs.slide_height = self._get_slide_dimensions_pptx(aspect_ratio)

            # Add each screenshot as a slide
            for idx, screenshot_bytes in enumerate(screenshots):
                logger.info(f"Adding slide {idx + 1}/{len(screenshots)} to PPTX")

                # Create blank slide layout
                blank_slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(blank_slide_layout)

                # Convert screenshot bytes to PIL Image
                img = Image.open(io.BytesIO(screenshot_bytes))

                # Optimize image based on quality setting
                img = self._optimize_image(img, quality)

                # Convert back to bytes
                img_buffer = io.BytesIO()
                img.save(img_buffer, format='PNG', optimize=True)
                img_buffer.seek(0)

                # Add image to slide (fill entire slide)
                slide.shapes.add_picture(
                    img_buffer,
                    left=0,
                    top=0,
                    width=prs.slide_width,
                    height=prs.slide_height
                )

                logger.info(f"Slide {idx + 1} added successfully")

            # Save to BytesIO
            pptx_buffer = io.BytesIO()
            prs.save(pptx_buffer)
            pptx_bytes = pptx_buffer.getvalue()

            # Save to file if path provided
            if output_path:
                output_path.write_bytes(pptx_bytes)
                logger.info(f"PPTX saved to: {output_path}")

            logger.info(f"PPTX generated successfully ({len(pptx_bytes)} bytes)")
            return pptx_bytes

        except Exception as e:
            logger.error(f"PPTX generation failed: {e}", exc_info=True)
            raise RuntimeError(f"PPTX generation failed: {e}") from e

    def _get_slide_dimensions_pptx(self, aspect_ratio: str) -> tuple[int, int]:
        """
        Get PowerPoint slide dimensions based on aspect ratio.

        Args:
            aspect_ratio: '16:9' or '4:3'

        Returns:
            Tuple of (width, height) in EMUs (English Metric Units)
        """
        if aspect_ratio == "16:9":
            return (self.PPTX_WIDTH, self.PPTX_HEIGHT)
        elif aspect_ratio == "4:3":
            # 4:3 aspect ratio
            width = Inches(10)
            height = Inches(7.5)  # 10 * 3/4
            return (width, height)
        else:
            # Default to 16:9
            return (self.PPTX_WIDTH, self.PPTX_HEIGHT)

    def _optimize_image(self, img: Image.Image, quality: str) -> Image.Image:
        """
        Optimize image based on quality setting.

        Args:
            img: PIL Image to optimize
            quality: Quality level - 'high', 'medium', or 'low'

        Returns:
            Optimized PIL Image
        """
        # Quality settings determine max dimensions
        quality_settings = {
            "high": 1.0,      # Full resolution (1920x1080)
            "medium": 0.75,   # 75% resolution (1440x810)
            "low": 0.5        # 50% resolution (960x540)
        }

        scale = quality_settings.get(quality, 1.0)

        if scale < 1.0:
            new_width = int(img.width * scale)
            new_height = int(img.height * scale)

            # Use high-quality downsampling
            img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
            logger.info(f"Image optimized to {new_width}x{new_height} ({quality} quality)")

        return img

    async def add_speaker_notes(
        self,
        presentation_id: str,
        slide_count: int,
        notes: list[str]
    ) -> bytes:
        """
        Generate PPTX with speaker notes.

        This is an enhanced version that adds speaker notes to slides.

        Args:
            presentation_id: The UUID of the presentation
            slide_count: Number of slides in the presentation
            notes: List of speaker notes for each slide

        Returns:
            PPTX file as bytes with speaker notes
        """
        logger.info("Generating PPTX with speaker notes...")

        # First generate base PPTX
        pptx_bytes = await self.generate_pptx(presentation_id, slide_count)

        # Load PPTX for modification
        prs = Presentation(io.BytesIO(pptx_bytes))

        # Add notes to each slide
        for idx, slide in enumerate(prs.slides):
            if idx < len(notes) and notes[idx]:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes[idx]
                logger.info(f"Added notes to slide {idx + 1}")

        # Save modified PPTX
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_bytes = pptx_buffer.getvalue()

        logger.info("Speaker notes added successfully")
        return pptx_bytes
